"""Generate the component & trust map deck: what the system is made of, which
upstream projects it leans on (ACRO/SACRO, OpenSAFELY, GOV.UK DS, Alloy, Lean,
Tailscale, the runtime five), and how far each part is trusted.

Same contract as make_decks.py: the .pptx output is gitignored and distributed
via releases; this generator is the source of truth. No screenshots needed, so
there is no capture step. Evidence numbers (red-team corpus size) are read from
the repository at build time.

    uv run --group decks python scripts/make_component_map.py     # ./artifacts
"""

from __future__ import annotations

import argparse
import os
import sys

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from make_decks import (
    BLUE, GREEN, GREY, INK, PANEL, SLIDE_H, SLIDE_W, WHITE,
    _bar, _bg, _text, table_slide, title_slide,
)

# A harness must never write to the operator's real audit log. `safetre_web.app`
# opens `SAFETRE_AUDIT_DB` at import and now appends a policy record there
# (#55), so merely importing it from a script pollutes `./audit.db` — which is
# hardening #36 all over again, and did happen (#57). Pin a throwaway path
# BEFORE the import, exactly as `tests/conftest.py` does for the test suite.
import os as _os          # noqa: E402
import tempfile as _tempfile  # noqa: E402

_os.environ.setdefault(
    "SAFETRE_AUDIT_DB",
    _os.path.join(_tempfile.gettempdir(), "safetre-harness-audit.db"))

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

RED = RGBColor(0xD4, 0x35, 0x1C)
GREEN_TINT = RGBColor(0xE7, 0xF5, 0xEC)
RED_TINT = RGBColor(0xFD, 0xEA, 0xEA)
BLUE_TINT = RGBColor(0xEE, 0xF6, 0xFF)


def _attack_counts() -> tuple[int, int]:
    entries = yaml.safe_load(open(os.path.join(ROOT, "redteam", "attacks.yaml")))
    benign = sum(1 for a in entries if a.get("type") == "benign")
    return len(entries) - benign, benign


# --- drawing helpers ------------------------------------------------------------

def _dash(line_format):
    ln = line_format._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))


def _box(slide, x, y, w, h, title, note, fill=WHITE, line=GREY, dash=False,
         title_size=11.5, note_size=8):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        shp.adjustments[0] = 0.06
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(1.25)
    if dash:
        _dash(shp.line)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = Pt(5)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(3)
    r = tf.paragraphs[0].add_run()
    r.text = title
    r.font.size = Pt(title_size)
    r.font.bold = True
    r.font.name = "Arial"
    r.font.color.rgb = INK
    if note:
        p = tf.add_paragraph()
        p.space_before = Pt(2)
        r = p.add_run()
        r.text = note
        r.font.size = Pt(note_size)
        r.font.name = "Arial"
        r.font.color.rgb = GREY
    return shp


def _panel(slide, x, y, w, h, fill, line, label, label_colour):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(1.5)
    shp.shadow.inherit = False
    _text(slide, x + Inches(0.12), y + Inches(0.05), w - Inches(0.24), Inches(0.3),
          [(label, 9.5, label_colour, True)])
    return shp


def _arrow(slide, x1, y1, x2, y2, colour=GREY, width=1.5, dash=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = colour
    conn.line.width = Pt(width)
    conn.shadow.inherit = False
    ln = conn.line._get_or_add_ln()
    if dash:
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    return conn


def _header(prs, title, subtitle, accent=BLUE):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _bar(s, accent)
    _text(s, Inches(0.6), Inches(0.32), Inches(12.2), Inches(0.6), [(title, 26, INK, True)])
    _text(s, Inches(0.6), Inches(0.95), Inches(12.2), Inches(0.4), [(subtitle, 13, GREY, False)])
    return s


def _legend(slide, y, items):
    x = Inches(0.6)
    for fill, line, dash, label in items:
        chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(0.3), Inches(0.18))
        chip.fill.solid()
        chip.fill.fore_color.rgb = fill
        chip.line.color.rgb = line
        chip.line.width = Pt(1.25)
        chip.shadow.inherit = False
        if dash:
            _dash(chip.line)
        _text(slide, x + Inches(0.38), y - Inches(0.05), Inches(2.9), Inches(0.3),
              [(label, 9.5, INK, False)])
        x += Inches(3.2)


# --- slide 2: the runtime pipeline ----------------------------------------------

def pipeline_slide(prs):
    s = _header(prs, "Runtime components — the safe-outputs pipeline",
                "Every box a request crosses; colour = how far it is trusted")

    _panel(s, Inches(1.85), Inches(1.42), Inches(9.6), Inches(4.78), BLUE_TINT, BLUE,
           "Trusted Research Environment — safepod boundary: no row-level egress · "
           "least-privilege systemd unit", BLUE)

    ra_y, ra_h = Inches(1.95), Inches(1.3)
    _box(s, Inches(0.25), ra_y, Inches(1.4), ra_h, "Researcher",
         "Safe People allowlist; asks in natural language.")
    _box(s, Inches(2.1), ra_y, Inches(1.65), ra_h, "Tailscale serve — restricted channel",
         "Real peer address checked against allowlist CIDRs; forwarded headers ignored.",
         fill=WHITE, line=BLUE)
    _box(s, Inches(3.95), ra_y, Inches(1.65), ra_h, "safetre_web — FastAPI + Jinja2",
         "GOV.UK-styled shell; CSP script-src 'self', no CDN. Outside the claim; frozen.",
         fill=WHITE, line=BLUE)
    _box(s, Inches(5.8), ra_y, Inches(1.7), ra_h, "LLM planner (local, unnamed)",
         "Untrusted by design: proposes QuerySpec JSON only — no code, SQL, files, "
         "sockets. Endpoint host allowlisted.", fill=RED_TINT, line=RED)
    _box(s, Inches(7.7), ra_y, Inches(1.75), ra_h, "QuerySpec validation — Pydantic v2",
         "The security boundary: extra='forbid'; anything off the catalogue is "
         "rejected before it runs.", fill=GREEN_TINT, line=GREEN)
    _box(s, Inches(9.65), ra_y, Inches(1.6), ra_h, "Legacy code path (exec)",
         "Red-team comparison only — never the web path; out of the claim (spec N3).",
         fill=WHITE, line=RED, dash=True)

    mid_a = ra_y + ra_h // 2
    for x1, x2 in ((1.65, 2.1), (3.75, 3.95), (5.6, 5.8), (7.5, 7.7)):
        _arrow(s, Inches(x1), mid_a, Inches(x2), mid_a)

    rb_y, rb_h = Inches(3.65), Inches(1.3)
    _box(s, Inches(7.7), rb_y, Inches(1.75), rb_h, "Query engine — DuckDB (read-only)",
         "Parameterised SQL over views that expose no identifying columns.",
         fill=GREEN_TINT, line=GREEN)
    _box(s, Inches(5.8), rb_y, Inches(1.7), rb_h, "Procedure registry — GLM · ANOVA",
         "Registered contracts (R14); models fit only on gateway-finalized cell "
         "tables (P19–P22); stdlib IRLS.", fill=GREEN_TINT, line=GREEN)
    _box(s, Inches(3.95), rb_y, Inches(1.65), rb_h, "Disclosure gateway",
         "Min cell size, dominance, rounding, suppression to fixpoint. "
         "ACRO-inspired stand-in — and says so.", fill=GREEN_TINT, line=GREEN)
    _box(s, Inches(2.1), rb_y, Inches(1.65), rb_h, "Session auditor",
         "Simulatable: decides from published marginals; differencing lineage + "
         "query budget.", fill=GREEN_TINT, line=GREEN)

    _arrow(s, Inches(8.575), ra_y + ra_h, Inches(8.575), rb_y)  # validation -> engine
    mid_b = rb_y + rb_h // 2
    for x1, x2 in ((7.7, 7.5), (5.8, 5.6), (3.95, 3.75)):
        _arrow(s, Inches(x1), mid_b, Inches(x2), mid_b)

    rc_y, rc_h = Inches(5.18), Inches(0.96)
    _box(s, Inches(2.1), rc_y, Inches(1.65), rc_h, "Human-in-the-loop",
         "Residual findings escalate; a denial renders no data (P18).",
         fill=GREEN_TINT, line=GREEN)
    _box(s, Inches(3.95), rc_y, Inches(1.65), rc_h, "ACRO / SACRO (DARE UK)",
         "Production SDC — roadmap item 1; replaces the stand-in.",
         fill=WHITE, line=GREY, dash=True)
    _box(s, Inches(5.8), rc_y, Inches(1.7), rc_h, "Audit log — HMAC-chained",
         "Every request, spec, decision; /api/audit/verify; key off-host.",
         fill=GREEN_TINT, line=GREEN)
    _box(s, Inches(7.7), rc_y, Inches(1.75), rc_h, "Row-level data — synthetic only",
         "make_data.py; pinned disclosure anchors; no real data yet.",
         fill=RED_TINT, line=RED)
    _box(s, Inches(0.25), rc_y, Inches(1.4), rc_h, "Released aggregate",
         "Only after every gate has passed.", fill=GREEN_TINT, line=GREEN)

    _arrow(s, Inches(2.925), rb_y + rb_h, Inches(2.925), rc_y)              # auditor -> human
    _arrow(s, Inches(4.775), rc_y, Inches(4.775), rb_y + rb_h, dash=True)   # ACRO -> gateway
    _arrow(s, Inches(8.575), rc_y, Inches(8.575), rb_y + rb_h, dash=True)   # data -> engine
    _text(s, Inches(8.65), Inches(4.96), Inches(1.2), Inches(0.25), [("read-only", 8, GREY, False)])
    _arrow(s, Inches(2.1), rc_y + rc_h // 2, Inches(1.65), rc_y + rc_h // 2,
           colour=GREEN, width=2)                                           # human -> released
    _arrow(s, Inches(0.95), rc_y, Inches(0.95), ra_y + ra_h, dash=True)     # back to researcher

    _legend(s, Inches(6.45), [
        (GREEN_TINT, GREEN, False, "deterministic control — the trust anchor"),
        (RED_TINT, RED, False, "untrusted or sensitive"),
        (WHITE, BLUE, False, "transport & shell — outside the claim"),
        (WHITE, GREY, True, "external / planned / red-team only"),
    ])
    _text(s, Inches(0.6), Inches(6.85), Inches(12.2), Inches(0.55),
          [("A literal QuerySpec (JSON request) skips the planner but crosses every "
            "control unchanged (R17). The planner model is deliberately unnamed in "
            "tracked content: naming it invites model-targeted prompting, and no "
            "safety property depends on the choice.", 9.5, GREY, False)])


# --- slide 3: the assurance toolchain --------------------------------------------

def assurance_slide(prs):
    n_attack, n_benign = _attack_counts()
    s = _header(prs, "Assurance toolchain — what checks the code",
                "All of this is dev/CI-side: none of it ships in the runtime install",
                accent=GREEN)

    _box(s, Inches(0.6), Inches(1.8), Inches(2.3), Inches(1.1), "safetre/ — research core",
         "5 runtime packages (pydantic, duckdb, pandas, numpy, pyyaml); uv.lock frozen.",
         fill=GREEN_TINT, line=GREEN)
    _box(s, Inches(0.6), Inches(3.35), Inches(2.3), Inches(1.3), "Generators",
         "gen_alloy_catalogue.py · gen_lean_catalogue.py export the finite request "
         "space from the live code; exit 1 on drift.")
    _arrow(s, Inches(1.75), Inches(2.9), Inches(1.75), Inches(3.35))

    _box(s, Inches(3.55), Inches(2.6), Inches(2.55), Inches(1.35), "Alloy Analyzer 6.2.0 — dev/CI",
         "Two bounded models: the GLM release path (P19 · P21 · P4) and the "
         "auditor's differencing rule with its documented residuals; jar fetched "
         "by pinned sha256.",
         fill=WHITE, line=INK)
    _box(s, Inches(3.55), Inches(4.35), Inches(2.55), Inches(1.45), "Lean 4 — pinned v4.32.0, dev/CI",
         "Query-boundary proofs: no identifiers (P3), internal columns never "
         "released (P4), single-view read-only SQL, all values bound (P9); "
         "414-case render-equality vs the live engine.",
         fill=WHITE, line=INK)
    _arrow(s, Inches(2.9), Inches(3.7), Inches(3.55), Inches(3.25))
    _arrow(s, Inches(2.9), Inches(4.3), Inches(3.55), Inches(4.9))

    _panel(s, Inches(6.55), Inches(1.75), Inches(6.45), Inches(4.6), PANEL, GREY,
           "GitHub Actions CI — actions SHA-pinned · contents:read token · fork PRs "
           "run with no secrets", INK)
    ci = [
        ("pytest + Hypothesis", "Unit, property-fuzz and exhaustive-enumeration "
         "suites; spec clauses cited in tests."),
        ("Red-team harness", f"{n_attack} attacks + {n_benign} benign baselines, "
         "gateway off vs on; exits nonzero on any failure."),
        ("Planner evaluation", "Scored corpus for the untrusted planner — refusal "
         "must come from the boundary."),
        ("Formal (Alloy + Lean)", "Model checks and lake build run next to pytest; "
         "toolchains fetched by pinned sha256; sync tests fail on drift."),
        ("bandit + pip-audit", "SAST plus CVE audit of the frozen dependency lock."),
        ("ruff", "Zero-findings lint baseline."),
        ("mkdocs build --strict", "Documentation links cannot rot."),
        ("pa11y", "WCAG 2.2 AA against the four demo states."),
    ]
    for i, (t, note) in enumerate(ci):
        x = Inches(6.75 + (i % 2) * 3.2)
        y = Inches(2.3 + (i // 2) * 1.0)
        _box(s, x, y, Inches(3.0), Inches(0.9), t, note)

    _text(s, Inches(0.6), Inches(6.55), Inches(12.2), Inches(0.6),
          [("Dev-only oracles (statsmodels · scipy) cross-validate the stdlib-only "
            "statistics; the runtime never imports them. Remaining formal work "
            "(roadmap item 2): the auditor's temporal model and value-level "
            "noninterference.", 9.5, GREY, False)])


# --- slides 4–5: upstream provenance ----------------------------------------------

def provenance_slides(prs):
    eco_rows = [
        ("ACRO / SACRO (DARE UK · AI-SDC)", "Production statistical disclosure control",
         "Planned",
         "Roadmap item 1: replaces the in-house stand-in gateway, slotting beneath "
         "the GLM layer. The stand-in states its own limits."),
        ("OpenSAFELY (Bennett Institute)", "TRE blueprint: code-to-data, checked outputs",
         "No — concept",
         "No code dependency; supplies the design assumptions this prototype "
         "re-tests with an AI analyst in the loop."),
        ("Five Safes (UK Data Service)", "Governance framework", "No — concept",
         "Vetting = Safe People/Projects; in-enclave model = Safe Settings; "
         "gateway = Safe Outputs."),
        ("GOV.UK Design System", "Interface idiom + accessibility bar", "Tokens only",
         "Hand-rolled CSS from the published specs (MIT); no govuk-frontend code, "
         "no GDS Transport or crown (restricted assets)."),
        ("Pydantic v2", "QuerySpec validation", "Yes",
         "The security boundary: extra='forbid', strict catalogue allowlists — "
         "anything else rejected before execution."),
        ("DuckDB", "Aggregate query engine", "Yes",
         "Opened read-only over views exposing no identifying columns; every filter "
         "value is a bound parameter."),
        ("pandas · numpy · PyYAML", "Data frames & config", "Yes",
         "Rest of the deliberately small 5-package runtime; frozen uv.lock, "
         "pip-audit in CI."),
        ("FastAPI · uvicorn · Jinja2 · httpx", "Demo web shell ('web' extra)", "Optional",
         "Reference deployment, not the security claim; CSP 'self', binds loopback, "
         "frozen — fixes only."),
    ]
    infra_rows = [
        ("Local LLM runtime (vLLM / llama.cpp class)", "Planner via OpenAI-compatible endpoint",
         "In enclave",
         "Untrusted by assumption; stdlib HTTP, no vendor SDK; endpoint host "
         "allowlisted; remote = opt-in, synthetic-only; model unnamed."),
        ("Tailscale", "Restricted channel + identity (reference deploy)", "Deploy",
         "App checks the real peer address against allowlist CIDRs and ignores "
         "forwarded headers; identity header trusted only behind tailscale serve."),
        ("Alloy Analyzer 6.2.0", "Bounded model checks (R16)", "Dev/CI",
         "Two models: GLM release path and the auditor's differencing rule; jar "
         "fetched in CI by pinned sha256; sync tests fail on drift."),
        ("Lean 4 (v4.32.0)", "Query-boundary proofs (P3, P4, P9)", "Dev/CI",
         "Toolchain sha256-pinned in CI; artifacts generated from live code; "
         "414-case render-equality against the engine; sync tests fail on drift."),
        ("GitHub Actions", "Continuous integration", "Dev/CI",
         "Actions pinned by commit SHA; least-privilege contents:read token; fork "
         "PRs run without secrets."),
    ]
    headers = ["Component", "Role", "In runtime?", "Trust & security notes"]
    widths = [2.6, 2.5, 1.0, 5.4]
    table_slide(prs, "Upstream projects 1/2 — ecosystem & runtime",
                headers, eco_rows, accent=BLUE, col_widths=widths)
    table_slide(prs, "Upstream projects 2/2 — model & assurance",
                headers, infra_rows, accent=BLUE, col_widths=widths)


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--out", default=os.path.join(ROOT, "artifacts"))
    args = ap.parse_args()
    os.makedirs(args.out, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(prs, "safe-tre-agent — components & trust map",
                "What the system is made of, which upstream projects it leans on, "
                "and how far each part is trusted")
    pipeline_slide(prs)
    assurance_slide(prs)
    provenance_slides(prs)

    out = os.path.join(args.out, "safe-tre-agent-components.pptx")
    prs.save(out)
    print(f"deck -> {out}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
