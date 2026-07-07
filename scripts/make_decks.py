"""Generate the presentation decks from current project state and screenshots.

Reproducible replacement for the hand-made binary decks: the .pptx outputs are
gitignored and distributed via releases, this generator is the source of truth.
Screenshots are captured live from the running web app.

    uv run --group decks python scripts/make_decks.py            # ./artifacts
    uv run --group decks python scripts/make_decks.py --shots-only

Regenerate screenshots against a running server (scripts/restart_web.sh), or
pass --no-capture to reuse existing PNGs in the shots directory.
"""

from __future__ import annotations

import argparse
import os
import subprocess
import sys
import urllib.request

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

# GOV.UK-derived palette, to match the interface the decks now show.
INK = RGBColor(0x0B, 0x0C, 0x0C)
BLUE = RGBColor(0x1D, 0x70, 0xB8)
GREEN = RGBColor(0x00, 0x70, 0x3C)
GREY = RGBColor(0x50, 0x5A, 0x5F)
PANEL = RGBColor(0xF3, 0xF2, 0xF1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

SHOTS = {
    "home": "/",
    "released": "/#q=mean%20spend%20by%20age%20band",
    "redacted": "/#q=mean%20spend%20by%20region%20and%20device%20os",
    "denied": "/#q=show%20mean%20wellbeing%20per%20donor",
}


def capture(shots_dir: str, base: str) -> None:
    """Screenshot each demo state with headless Chrome."""
    os.makedirs(shots_dir, exist_ok=True)
    chrome = next((p for p in ("/usr/bin/google-chrome", "/usr/bin/chromium",
                               "/usr/bin/chromium-browser") if os.path.exists(p)), None)
    if chrome is None:
        sys.exit("no chrome/chromium found for screenshots; pass --no-capture")
    try:
        urllib.request.urlopen(base + "/healthz", timeout=3)
    except Exception:
        sys.exit(f"web app not reachable at {base}; start it with scripts/restart_web.sh")
    for name, path in SHOTS.items():
        subprocess.run(
            [chrome, "--headless=new", "--disable-gpu", "--hide-scrollbars",
             "--window-size=1280,1400", "--run-all-compositor-stages-before-draw",
             "--virtual-time-budget=9000",
             f"--screenshot={os.path.join(shots_dir, name + '.png')}", base + path],
            check=True, stderr=subprocess.DEVNULL)
    print(f"screenshots -> {shots_dir}")


# --- slide helpers ------------------------------------------------------------

def _bg(slide, colour):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = colour


def _text(slide, left, top, width, height, runs, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, size, colour, bold) in enumerate(runs):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = colour
        run.font.bold = bold
        run.font.name = "Arial"
    return box


def _bar(slide, colour, height=Inches(0.18)):
    shp = slide.shapes.add_shape(1, 0, 0, SLIDE_W, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = colour
    shp.line.fill.background()


def _image_fit(slide, path, left, top, max_w, max_h):
    from PIL import Image
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(max_w / iw, max_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    slide.shapes.add_picture(path, left + (max_w - w) // 2, top, width=w, height=h)


def title_slide(prs, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, INK)
    _bar(s, BLUE, Inches(0.25))
    _text(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(2),
          [(title, 40, WHITE, True)])
    _text(s, Inches(0.9), Inches(4.2), Inches(11.5), Inches(1.5),
          [(subtitle, 20, RGBColor(0x99, 0xF6, 0xE4), False)])
    _text(s, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.6),
          [("Research prototype · synthetic data · not a government service", 12, GREY, False)])


def bullets_slide(prs, title, bullets, accent=BLUE):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _bar(s, accent)
    _text(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(1),
          [(title, 30, INK, True)])
    runs = [(("•  " + b), 18, INK, False) for b in bullets]
    box = _text(s, Inches(0.9), Inches(1.8), Inches(11.5), Inches(5), runs)
    for p in box.text_frame.paragraphs:
        p.space_after = Pt(12)


def shot_slide(prs, title, caption, image, accent=BLUE):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _bar(s, accent)
    _text(s, Inches(0.9), Inches(0.45), Inches(11.5), Inches(0.8),
          [(title, 26, INK, True)])
    _text(s, Inches(0.9), Inches(1.15), Inches(11.5), Inches(0.6),
          [(caption, 15, GREY, False)])
    if os.path.exists(image):
        _image_fit(s, image, Inches(0.9), Inches(1.8), Inches(11.5), Inches(5.3))
    else:
        _text(s, Inches(0.9), Inches(3), Inches(11.5), Inches(1),
              [(f"[missing screenshot: {os.path.basename(image)}]", 16, GREY, False)])


def table_slide(prs, title, headers, rows, accent=BLUE, col_widths=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _bar(s, accent)
    _text(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(1),
          [(title, 30, INK, True)])
    n_rows, n_cols = len(rows) + 1, len(headers)
    gt = s.shapes.add_table(n_rows, n_cols, Inches(0.9), Inches(1.7),
                            Inches(11.5), Inches(0.4 * n_rows)).table
    if col_widths:
        for c, w in enumerate(col_widths):
            gt.columns[c].width = Inches(w)
    for c, h in enumerate(headers):
        cell = gt.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = accent
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = h
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = WHITE
        r.font.name = "Arial"
    for ri, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = gt.cell(ri, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 else PANEL
            p = cell.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = val
            r.font.size = Pt(12)
            r.font.color.rgb = INK
            r.font.name = "Arial"


# --- decks --------------------------------------------------------------------

def build_overview(shots: str, out: str) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(prs, "Safe outputs gateway",
                "A safe-outputs gateway for automated analysis inside a Trusted Research Environment")

    bullets_slide(prs, "The problem", [
        "TRE disclosure control (min cell size, suppression) assumes a human analyst.",
        "An automated planner between analyst and data adds attack surface: injection carried in the data,",
        "   multi-query differencing, and code that smuggles rows into an ‘aggregate’.",
        "Question: does automation break the disclosure guarantee — and can a gateway restore it?",
    ])

    bullets_slide(prs, "The design", [
        "The planner only proposes a typed QuerySpec over an allowlisted catalogue — no code, no SQL.",
        "Validation rejects anything off-allowlist before execution; read-only DuckDB, bound parameters.",
        "Safe-outputs gateway: min donor count, dominance, influence, suppression — fail closed.",
        "Session auditor flags differencing from published (simulatable) marginals; refusals carry no numbers.",
        "Every request written to an HMAC-chained audit log.",
    ], accent=GREEN)

    shot_slide(prs, "The interface", "GOV.UK Design System, unbranded · WCAG 2.2 AA (pa11y-clean)",
               os.path.join(shots, "home.png"))
    shot_slide(prs, "A released query",
               "The gateway checks all pass; the answer is released with a magnitude table.",
               os.path.join(shots, "released.png"), accent=GREEN)
    shot_slide(prs, "A redacted query",
               "Small cells suppressed and a margin protected; released with a confidentiality note.",
               os.path.join(shots, "redacted.png"), accent=RGBColor(0xB1, 0x8A, 0x00))
    shot_slide(prs, "A denied query",
               "‘Wellbeing per donor’ is rejected at validation; no data table is rendered.",
               os.path.join(shots, "denied.png"), accent=RGBColor(0xD4, 0x35, 0x1C))

    bullets_slide(prs, "Statistical models, same guarantee (new)", [
        "GLMs (gaussian / logistic / Poisson) now run behind the gateway as registered procedures.",
        "Cells-first: a model is fitted ONLY from disclosure-vetted cell tables — never rows.",
        "If any design cell would be suppressed, the whole model is refused, loudly.",
        "Every release ships the vetted cell table it was fitted from: the analyst can reproduce it.",
        "One-way ANOVA just joined as a second tool — reusing the GLM's vetted cells,",
        "   adding only its own arithmetic: evidence the extension seam works.",
    ], accent=GREEN)

    bullets_slide(prs, "Evidence", [
        "Red-team: 17/17 attacks blocked with the gateway on; 7/20 leak row-level data with it off.",
        "Specification: 16 requirements, 22 prohibitions, each traced to code and a test.",
        "Planner evaluation: the planner proposes usefully but rarely refuses —",
        "   so refusal must come from the boundary, not the planner.",
        "360+ tests, red-team, an Alloy model check, strict docs build and pa11y all run in CI.",
    ])

    bullets_slide(prs, "What’s next", [
        "1. Integrate ACRO for production-grade statistical disclosure control.",
        "2. Extend the machine-checked model from the GLM path to the full query boundary.",
        "3. A differential-privacy accountant to close the simulatability residual.",
        "4. Cross-session and cross-user lineage.",
    ], accent=GREEN)

    prs.save(out)
    print(f"deck -> {out}")


AMBER = RGBColor(0xB1, 0x8A, 0x00)
RED = RGBColor(0xD4, 0x35, 0x1C)


def build_technical(shots: str, out: str) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(prs, "Safe outputs gateway — technical",
                "The boundary, the disclosure controls, and how they are verified")

    bullets_slide(prs, "Request lifecycle", [
        "Natural-language request → intent vetting (defence in depth).",
        "Planner (untrusted, automated) proposes a QuerySpec — the only executable output.",
        "Validation: Pydantic allowlist, extra=forbid — reject before running anything.",
        "Engine: validated spec → parameterised, read-only DuckDB over public views.",
        "Safe-outputs gateway → session auditor → human-in-the-loop → HMAC-chained log.",
    ])

    table_slide(prs, "Threat model (selected)",
                ["#", "Threat", "Control"],
                [["1", "Arbitrary code / RCE", "the planner writes no code; only a typed QuerySpec"],
                 ["2", "SQL injection", "bound parameters; identifiers regex-checked"],
                 ["3", "Identifier / free-text egress", "absent from every allowlist and view"],
                 ["4", "Small-cell / dominance", "min donor count, p%-rule, influence bound"],
                 ["5", "Differencing / triangulation", "simulatable session auditor + budget"],
                 ["6", "Injection via hostile data", "the planner can only emit a QuerySpec"],
                 ["9", "Tamper with the audit record", "HMAC-keyed chain, off-box anchor"],
                 ["17", "Fail-open suppression", "unresolved check → +inf → suppressed"]],
                col_widths=[0.7, 4.3, 6.5])

    bullets_slide(prs, "The QuerySpec boundary", [
        "A finite catalogue: 3 datasets, typed dimensions and measures, bounded group-by/filters.",
        "Registered procedures only: count, mean, sum, sum of squares, Pearson correlation.",
        "Direct identifiers, free text and raw timestamps are in no allowlist, so no valid query names them.",
        "The query space is finite and enumerable — which is what makes it testable and provable.",
    ], accent=GREEN)

    bullets_slide(prs, "The disclosure gateway", [
        "Minimum cell size counted over distinct individuals, not rows.",
        "Dominance (p%-rule) for sums and means; leave-one-out influence for correlations.",
        "Primary and complementary suppression so a margin cannot reconstruct a suppressed cell.",
        "Fail closed: an unresolved safety statistic is treated as unsafe and suppressed.",
    ], accent=GREEN)

    bullets_slide(prs, "Models behind the same gateway (GLM + one-way ANOVA)", [
        "GLMs (gaussian / logistic / Poisson over categorical terms) fit ONLY on gateway-vetted cell tables.",
        "Any suppressed design cell denies the whole model — no merging, no dropping, no silent repair.",
        "A release carries the vetted cell table: refitting from it reproduces the coefficients bit-for-bit.",
        "So the disclosure claim is inherited from the gateway — not re-argued per statistic.",
        "One-way ANOVA is a second tool on the same seam: same vetted mean / sum-of-squares cells,",
        "   one new stdlib numeric (the F-tail, cross-validated vs scipy), engine/gateway untouched.",
        "Machine-checked: exhaustive 767-point model skeleton (718 GLM + 49 ANOVA), refit-equality, Alloy in CI.",
    ], accent=GREEN)

    bullets_slide(prs, "Simulatable session auditing", [
        "Each released cohort is remembered by its normalised filter predicate.",
        "A new cohort within a small symmetric difference of a prior one is denied (differencing).",
        "The decision uses only published donor marginals — reproducible by the analyst.",
        "Refusals carry no numbers; the residual (sub-threshold isolation) is what DP closes.",
    ], accent=GREEN)

    shot_slide(prs, "The pipeline, made legible",
               "Each gateway stage reports a text status; a denial stops the request and renders no data.",
               os.path.join(shots, "denied.png"), accent=RED)

    bullets_slide(prs, "Verification", [
        "Normative specification: 16 requirements, 22 prohibitions, each traced to code and a test.",
        "Property-based tests sample the query space; exhaustive enumeration checks both skeletons.",
        "Red-team harness replays 20 scenarios, gateway off vs on — a CI gate.",
        "A bounded Alloy model (generated from the committed skeleton) checks P19/P21/P4 in CI.",
        "Strict docs build and pa11y (WCAG 2.2 AA) also run in CI.",
    ])

    prs.save(out)
    print(f"deck -> {out}")


def build_best_practice(shots: str, out: str) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(prs, "Safe outputs gateway — best practice",
                "Conformance with TRE security guidance")

    table_slide(prs, "Mapped to the Five Safes",
                ["Safe", "In this system"],
                [["Safe Projects", "intent vetting rejects blocked purposes pre-planning"],
                 ["Safe People", "identity allowlist behind the restricted channel"],
                 ["Safe Settings", "planner runs inside the safepod; read-only engine, no egress"],
                 ["Safe Data", "synthetic; identifiers and free text never queryable"],
                 ["Safe Outputs", "disclosure gateway + session auditor + human review"]],
                col_widths=[3.0, 8.5])

    bullets_slide(prs, "Where it already follows best practice", [
        "The untrusted-planner boundary is the published Action-Selector pattern (Beurer-Kellner 2025).",
        "The posture matches published OWASP guidance for untrusted automated components:"
        "   validation, least privilege, output checks, red-team.",
        "The SDC rules mirror the ACRO check set and the SDC Handbook.",
        "A frequency threshold of 10 is at or above the handbook's 3–5 and OpenSAFELY's redact-≤7.",
        "Governance is honest: synthetic-only, HMAC-chained audit, stated limitations.",
    ], accent=GREEN)

    table_slide(prs, "Known deviations, stated openly",
                ["#", "Deviation", "Status"],
                [["D1", "Auditor simulatability", "fixed — decides from published marginals"],
                 ["D2", "Cumulative disclosure bound", "roadmap — differential-privacy accountant"],
                 ["D5", "One checker vs two humans", "human-in-the-loop present; reviewer queue planned"],
                 ["D6", "Influence threshold unvalidated", "documented; ACRO integration supersedes"]],
                col_widths=[0.7, 5.3, 5.5], accent=AMBER)

    bullets_slide(prs, "Grounded in existing infrastructure", [
        "OpenSAFELY (Bennett Institute, Oxford) — the code-to-data, outputs-checked TRE model.",
        "ACRO / SACRO (DARE UK) — the target production-grade disclosure-control dependency.",
        "Five Safes — the governance framing.",
        "This project adds the agent-aware layer above that established practice.",
    ])

    bullets_slide(prs, "What not to overclaim", [
        "These are statistical disclosure controls, not differential privacy — no epsilon budget yet.",
        "The auditor does not defend across sessions or colluding users.",
        "The disclosure engine is a stand-in for ACRO.",
        "The legacy code-execution path is a red-team narrative, not a secure sandbox.",
    ], accent=RED)

    prs.save(out)
    print(f"deck -> {out}")


def build_elif(shots: str, out: str) -> None:
    """The plain-language deck: docs/elif.md as slides (TL;DR + ELI5)."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(prs, "What we built, explained simply",
                "Statistical models behind the safe-outputs gateway — the ELIF version "
                "(precise version: docs/specification.md)")

    bullets_slide(prs, "TL;DR", [
        "The gateway now fits statistical models (gaussian / logistic / Poisson GLMs).",
        "A model is computed ONLY from summary cells the gateway already checked and released.",
        "If even one cell would be blocked, the whole model is refused — loudly, never quietly repaired.",
        "Every release ships the cell table it was fitted from: refitting reproduces it bit-for-bit.",
        "Around it: a framework where every statistical procedure is a registered, CI-enforced contract.",
    ], accent=GREEN)

    bullets_slide(prs, "The library with a careful librarian", [
        "A library holds everyone's diaries. You may never read a diary.",
        "You may ask about GROUPS — and the librarian follows strict rules:",
        "   no answers about groups smaller than ten; answers rounded; every question remembered.",
        "An automated helper turns your English into her official request form.",
        "We do NOT trust the helper — it can only fill in the form; the rulebook checks every box.",
    ])

    bullets_slide(prs, "The trick that makes models safe", [
        "Fitting a model normally means reading every individual's data.",
        "For the models we allow, the maths has a special property:",
        "   the model is computable EXACTLY from the group summaries alone (verified to ~1e-14).",
        "So the fit uses only tables the librarian already released — it cannot know more",
        "than she already said out loud. That is the whole safety argument.",
    ], accent=GREEN)

    bullets_slide(prs, "How a model request runs", [
        "1. Work out which group tables the model needs.",
        "2. Ask for each through EXACTLY the same rulebook as any hand query.",
        "3. Any blocked cell ⇒ refuse the whole model, out loud (no merging, no dropping).",
        "4. Fit with a pure calculator that physically cannot touch the database.",
        "5. Release coefficients + the very table they came from.",
    ])

    shot_slide(prs, "What the researcher sees",
               "The same gateway page: models appear as a released result with "
               "every check reported; a refusal renders no data at all.",
               os.path.join(shots, "released.png"), accent=GREEN)

    bullets_slide(prs, "Why refusing loudly matters", [
        "A 'helpful' system might quietly drop the too-small group and answer anyway.",
        "Then you'd trust an answer to a question you didn't ask.",
        "Here: blocked means blocked, and it says so.",
        "Even the refusal is computed only from things you were allowed to know —",
        "so a 'no' can't leak a secret either.",
    ], accent=AMBER)

    bullets_slide(prs, "The framework: adding statistics without adding holes", [
        "Every procedure is a registered contract: allowed columns, blessed query shape,",
        "safety witnesses (or inheritance), declared outputs, and its finite request space.",
        "Skip an obligation and the BUILD FAILS — a test enumerates and demands each one.",
        "Then: all 767 model shapes tried exhaustively · random fuzzing · 20 replayed attacks ·",
        "an Alloy solver searches for any way to fit past a blocked cell (it finds none —",
        "and finds the counterexample instantly when we deliberately weaken the rule).",
    ], accent=GREEN)

    bullets_slide(prs, "The framework in action: adding one-way ANOVA", [
        "ANOVA asks: do the averages differ across groups? (e.g. does spend differ by age band?)",
        "We added it WITHOUT touching the safe core — because it needs only the group summaries",
        "   (each group's average, spread and size) the librarian ALREADY checks and releases.",
        "So genuinely new was tiny: one formula for the score, one list of which output is which.",
        "It inherits every rule for free: one too-small group ⇒ the whole test is refused, out loud;",
        "   and you can redo the test yourself from the table it hands back.",
        "That is the seam working: a new statistic, no new way to leak.",
    ], accent=GREEN)

    bullets_slide(prs, "And one embarrassing thing we found", [
        "The existing counting rule rounded the count in one column…",
        "…and wrote the EXACT count in the column next to it.",
        "Count rounding was doing nothing. Found while planning, fixed first, regression-tested.",
        "Exactly the gap the new 'declare every released column' contract makes impossible.",
    ], accent=RED)

    table_slide(prs, "The numbers",
                ["What", "Count"],
                [["Spec clauses", "R1–R16, P1–P22 (7 new this round)"],
                 ["Model shapes, all machine-checked", "767 (718 GLM + 49 ANOVA)"],
                 ["Tests in the default suite", "360+ (plus exhaustive -m slow pass)"],
                 ["Red-team attacks, all blocked by a named control", "20 (9 new)"],
                 ["Solver-checked properties (Alloy, in CI)", "4"],
                 ["Agreement with reference implementations", "~1e-14 exact / 1e-8 tested"],
                 ["New runtime dependencies", "0 — the fitter is stdlib-only"]],
                col_widths=[7.0, 4.5])

    bullets_slide(prs, "What it still does not do", [
        "Logistic/Poisson models with continuous predictors (genuinely need rows — parked for ACRO).",
        "Cross-session or colluding-user protection (differential privacy is the roadmap answer).",
        "It remains a research prototype, on synthetic data only.",
    ], accent=GREY)

    prs.save(out)
    print(f"deck -> {out}")


def build_guidelines(shots: str, out: str) -> None:
    """The maintainer's deck: docs/maintenance.md as slides."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(prs, "Maintenance guidelines",
                "How to change safe-tre-agent without breaking its claim — "
                "the playbook (precise version: docs/maintenance.md)")

    bullets_slide(prs, "Two principles govern every change", [
        "The SPECIFICATION leads: behaviour changes start as a clause amendment with a",
        "   traceability row; code and tests cite the clause; prose docs follow.",
        "Refuse loudly, never repair silently: a merged category, a dropped cell, or a",
        "   substituted question is wrong even when it is 'safe'. Denials are part of the product.",
        "Boundary files (query/engine/procedures/disclosure/audit/identity) always get human review.",
    ])

    table_slide(prs, "The gate list — green before every merge",
                ["Gate", "Command"],
                [["Lint (zero findings)", "uv run ruff check safetre safetre_web tests …"],
                 ["Tests + exhaustive skeletons", "uv run pytest -q  ·  uv run pytest -q -m slow"],
                 ["Red-team, gateway off vs on", "uv run python redteam/run_redteam.py"],
                 ["SAST + dependency CVEs", "uv run bandit -q -r …  ·  uv run pip-audit"],
                 ["Docs (strict)", "uv run --group docs mkdocs build --strict"],
                 ["Formal sync + model check", "gen_alloy_catalogue.py · formal/run_checks.py"],
                 ["Accessibility (UI changes)", "pa11y on home / released / redacted / denied"]],
                col_widths=[4.0, 7.5])

    bullets_slide(prs, "Adding a statistical procedure (1/2) — architecture first", [
        "Computable from catalogued aggregates?  Then CELLS-FIRST is mandatory:",
        "   plan ordinary QuerySpecs, fit only on gateway-finalized tables,",
        "   inherit every disclosure rule, prove reproducibility (refit-equality).",
        "Genuinely needs rows?  That route waits for ACRO (roadmap item 1).",
        "Spec first: amend R/P clauses + traceability BEFORE writing code.",
    ], accent=GREEN)

    bullets_slide(prs, "Adding a statistical procedure (2/2) — the checklist", [
        "Register the contract in procedures.py (fragments only — the SafeSQL shape is central).",
        "Declare obligations in test_procedure_conformance.py — double-entry: build fails on mismatch.",
        "Regenerate formal artifacts (gen_alloy_catalogue.py --write); sync tests gate drift.",
        "Extend EVERY layer: exhaustive enumeration · Hypothesis · oracle · red-team · eval corpus.",
        "Advertise last: manifest + planner prompt bump in ONE isolated, reviewed commit.",
        "Never: per-observation outputs · disclosure logic inside a procedure · silent fn fall-through.",
    ], accent=GREEN)

    bullets_slide(prs, "Worked example: one-way ANOVA (this release)", [
        "Test: computable from catalogued aggregates?  YES — it needs only per-group",
        "   mean, Σ(y²) and n: EXACTLY the gaussian GLM's design cells. So: cells-first, mandatory.",
        "Plan the same two design-cell QuerySpecs; inherit the gateway, P19/P21/P22, reproducibility.",
        "New code stayed small: stats.f_sf (F-tail, reusing the incomplete beta — no new dependency),",
        "   AnovaSpec (typed boundary), AnovaProcedure. engine / disclosure / service: UNTOUCHED.",
        "Checklist discharged: conformance obligation · skeleton regen (+49 pts) · oracle vs scipy.f_oneway ·",
        "   noninterference guard · manifest planned→available (v5). Written up: docs/adding-a-statistical-tool.md.",
    ], accent=GREEN)

    bullets_slide(prs, "Datasets, columns, thresholds", [
        "Every column gets a DI/QI/S/R role and (categorical) a declared domain.",
        "Identifiers, free text, raw timestamps: in NO allowlist, NO public view — ever.",
        "Update the public view AND the _u unit view; regenerate skeleton + Alloy block.",
        "Keep the synthetic disclosure anchors (sub-threshold NI / sex-X, hostile strings).",
        "Threshold floors are pinned by tests (min cell ≥ 5, rounding ≥ 5, dominance ≤ 0.5);",
        "   unresolved safety values fail CLOSED (+inf) — preserve the pattern.",
    ])

    bullets_slide(prs, "Manifest & planner prompt = the safepod contract", [
        "A tool is live only when it moves from planned_tool_classes into tools[].",
        "Treat any change as a mini release: one isolated commit, MANIFEST_VERSION bump,",
        "   contract pins updated deliberately, security review, signed artifact in a real safepod.",
        "Prompt changes carry no safety obligation (the planner is untrusted) — but a quality one:",
        "   run the planner eval before and after; extend the corpus for new phrasings.",
    ], accent=AMBER)

    bullets_slide(prs, "Changing the UI (e.g. GOV.UK guidance changes)", [
        "The shell is FROZEN — fixes only. Design-system updates, a11y fixes = fixes.",
        "   New controls, pages, or data surfaces = not fixes.",
        "Stay unbranded: GOV.UK layout, but no GDS Transport, no crown, no government claim.",
        "CSP stays script-src 'self' — no CDN assets, ever; identity/channel code is off-limits.",
        "A denial renders NO data table (P18). Gate on pa11y (WCAG 2.2 AA), all four demo states.",
        "Afterwards: regenerate deck screenshots (temporary allow-all identity, then restore).",
    ], accent=AMBER)

    bullets_slide(prs, "Install & deploy — the non-negotiables", [
        "Pinned env only: uv sync --all-extras --frozen. Runtime deps stay at five packages;",
        "   oracles and tooling are dev-only; actions and the Alloy jar are SHA-pinned.",
        "A real deployment runs the planner locally (a remote endpoint is an egress channel;",
        "   synthetic-only, explicit opt-in). Unreachable backend ⇒ fail loudly, never the mock.",
        "Identity fails closed: header trusted on loopback or behind an asserted proxy only.",
        "Audit key and chain anchor live OFF-BOX; least-privilege systemd unit; loopback bind.",
    ])

    bullets_slide(prs, "Releasing · and when something goes wrong", [
        "Release: CHANGELOG → version bump (single-sourced) → full gates → uv build →",
        "   install the wheel in a FRESH venv and smoke it → tag → DRAFT GitHub release.",
        "Publishing (and PyPI) is a deliberate human decision.",
        "Security finding: numbered hardening-log entry + a pinning regression test +",
        "   where possible a framework change that makes the whole class inexpressible",
        "   (hardening #25 → declared output contracts is the model). Never fix a leak silently.",
    ], accent=GREEN)

    prs.save(out)
    print(f"deck -> {out}")


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--out-dir", default=os.path.join(ROOT, "artifacts"))
    ap.add_argument("--shots-dir", default=os.path.join(ROOT, "artifacts", "shots"))
    ap.add_argument("--base-url", default="http://127.0.0.1:8800")
    ap.add_argument("--no-capture", action="store_true",
                    help="reuse existing screenshots instead of capturing")
    ap.add_argument("--shots-only", action="store_true")
    args = ap.parse_args()

    if not args.no_capture:
        capture(args.shots_dir, args.base_url)
    if args.shots_only:
        return 0

    os.makedirs(args.out_dir, exist_ok=True)
    build_overview(args.shots_dir,
                   os.path.join(args.out_dir, "safe-tre-agent-overview.pptx"))
    build_technical(args.shots_dir,
                    os.path.join(args.out_dir, "safe-tre-agent-technical.pptx"))
    build_best_practice(args.shots_dir,
                        os.path.join(args.out_dir, "safe-tre-agent-best-practice.pptx"))
    # the committed .ppt explainers (deliberately outside the artifacts/*.pptx
    # ignore rule): the plain-language deck alongside docs/elif.md, and the
    # maintainer's deck alongside docs/maintenance.md.
    build_elif(args.shots_dir, os.path.join(args.out_dir, "ELIF.ppt"))
    build_guidelines(args.shots_dir, os.path.join(args.out_dir, "GUIDELINES.ppt"))
    return 0


if __name__ == "__main__":
    sys.exit(main())
