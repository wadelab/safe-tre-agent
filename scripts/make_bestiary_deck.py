"""Build the plain-language bestiary deck: the security problems, as monsters.

`docs/bestiary.md` is the field guide — precise, indexed to the hardening log,
and long. This is the version you can give someone in twenty minutes: one
creature per slide, its card, what it wants in plain English, and the one
sentence about what stops it.

The audience is a research colleague, an IG officer or a funder — someone who
needs to understand *what kinds of thing go wrong with a data gateway* without
learning the vocabulary first. So the rule here is stricter than the field
guide's: **no term is used before it is explained, and every claim is one a
reader could check in the hardening log.** The metaphors are memory aids; the
parenthetical finding numbers are the audit trail.

Cards come from `docs/figures/bestiary/`, the committed derivatives written by
`scripts/make_bestiary_cards.py` — 600 px or better, which is enough at the
size a slide shows them.

    uv run --group decks python scripts/make_bestiary_deck.py
"""

from __future__ import annotations

import argparse
import os
import sys

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from make_decks import (  # noqa: E402
    BLUE, GREEN, GREY, INK, PANEL, SLIDE_H, SLIDE_W, WHITE,
    _bar, _bg, _text, bullets_slide, title_slide,
)

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
CARDS = os.path.join(ROOT, "docs", "figures", "bestiary")


def _fit_centred(slide, path: str, left, top, max_w, max_h) -> None:
    """Fit an image inside a box, centred both ways.

    The cards run from 0.71 (tall portrait) to 1.49 (landscape), so a fixed box
    and a top-aligned image strands the wide ones against the slide's ceiling
    with a hand's width of nothing beneath them. `make_decks._image_fit` centres
    horizontally only, which is right for screenshots — they are all one shape.
    """
    import io

    from PIL import Image

    # The committed .png are ~3 MB print masters; a slide needs nowhere near
    # that. Prefer the 600 px .webp derivative and hand python-pptx a small
    # in-memory JPEG (it cannot read webp), which keeps the deck emailable.
    stem = os.path.splitext(path)[0]
    webp = stem + ".webp"
    src = webp if os.path.exists(webp) else path
    with Image.open(src) as im:
        iw, ih = im.size
        buf = io.BytesIO()
        im.convert("RGB").save(buf, "JPEG", quality=88, optimize=True)
    buf.seek(0)
    scale = min(max_w / iw, max_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    slide.shapes.add_picture(buf, left + (max_w - w) // 2, top + (max_h - h) // 2,
                             width=w, height=h)


def card_slide(prs, card: str, name: str, wants: str, story: list[str],
               cage: str, accent=BLUE) -> None:
    """One creature: its card on the left, plain English on the right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE)
    _bar(slide, accent)

    art = os.path.join(CARDS, card)
    if os.path.exists(art):
        _fit_centred(slide, art, Inches(0.6), Inches(0.75), Inches(4.0), Inches(6.2))
    else:
        _text(slide, Inches(0.6), Inches(3.2), Inches(4.0), Inches(1),
              [(f"[missing card: {card}]", 13, GREY, False)])

    left, width = Inches(5.1), Inches(7.6)
    _text(slide, left, Inches(0.75), width, Inches(0.9), [(name, 32, INK, True)])
    _text(slide, left, Inches(1.6), width, Inches(0.7),
          [(f"What it wants: {wants}", 16, accent, True)])

    box = _text(slide, left, Inches(2.4), width, Inches(3.0),
                [(line, 17, INK, False) for line in story])
    for para in box.text_frame.paragraphs:
        para.space_after = Pt(10)

    panel = slide.shapes.add_shape(1, left, Inches(5.5), width, Inches(1.45))
    panel.fill.solid()
    panel.fill.fore_color.rgb = PANEL
    panel.line.fill.background()
    _text(slide, left + Inches(0.25), Inches(5.7), width - Inches(0.5), Inches(1.1),
          [("The cage", 13, GREY, True), (cage, 16, INK, False)])


# (card, name, wants, story, cage, accent)
CREATURES = [
    ("01_the_subtractor.png", "The Subtractor", "two answers that differ by one person",
     ["Ask how much everyone in a region spent. Then ask again, excluding",
      "the over-50s. Subtract one answer from the other and you have the",
      "spending of a handful of people — without ever asking about them.",
      "Every single question was reasonable. Only the subtraction was theft."],
     "The gateway remembers what it has already told you, and refuses an answer "
     "whose difference from an earlier one would isolate too few people (#40).",
     BLUE),

    ("02_the_whale.png", "The Whale", "to be personally visible inside a total",
     ["In a group of twenty, one person accounts for most of the spending.",
      "Publish the group total and you have very nearly published theirs.",
      "The group passed the 'at least ten people' rule, so counting alone",
      "never notices — you have to look at who is inside the number."],
     "A group is suppressed when one person is more than half of it — measured by "
     "size, so a large refund counts as much as a large payment (#41).",
     GREEN),

    ("03_the_masker.png", "The Masker", "to act while wearing someone else's name",
     ["Someone knocks and says 'I am Dr Smith'. If the badge is taken on",
      "trust, anyone can be Dr Smith — and the record of what they did will",
      "have Dr Smith's name on it.",
      "Worse: the limits on how much any one person may ask are counted",
      "per name. A new name is a fresh allowance."],
     "The badge must be countersigned by the doorkeeper, and sharing a corridor "
     "with the system is not proof of identity (#45).",
     INK),

    ("04_the_nixie.png", "The Nixie", "to make the system print a name",
     ["Sometimes the secret is not the number, it is the label beside it.",
      "'People aged 41 in this small county: 12' has already said something",
      "before you read the 12.",
      "A typo, or a hostile string typed into a form, is a label too."],
     "Only words from the official codebook may be printed as labels; anything "
     "else is withheld, however many people it covers (#43, #29).",
     BLUE),

    ("05_the_sphinx.png", "The Sphinx", "to talk its way past you",
     ["The system has an AI assistant that turns your English into a formal",
      "request. The Sphinx does not attack the rules — it asks the assistant",
      "nicely for the raw records, and hopes it obliges.",
      "You cannot fix this by making the assistant more sensible."],
     "The assistant can only tick boxes on a fixed form. There is no box for "
     "'raw records', so there is nothing for it to agree to.",
     GREEN),

    ("06_the_white_rabbit.png", "The White Rabbit", "to read the clock, not the answer",
     ["You can learn things from how long a reply takes. A question about a",
      "big group takes longer than one about a small group — so timing can",
      "put secret groups in size order, which is exactly what hiding their",
      "sizes was for."],
     "Every reply waits for the same clock tick before it is sent, and work that "
     "runs too long is refused at that same moment rather than when it finishes "
     "(#54).",
     GREY),

    ("07_the_ghost.png", "The Ghost", "to happen without being written down",
     ["Every request is meant to leave a line in a tamper-proof log. If some",
      "requests can slip through unlogged — a crash, say — then that is",
      "precisely where you would do the thing you did not want recorded.",
      "A log of almost everything has a hole exactly where it matters."],
     "Every request writes exactly one line, crashes included, and the line "
     "records the failure's type but never its message (#37).",
     INK),

    ("08_the_hydra.png", "The Hydra", "to survive the fix aimed at it",
     ["A problem is reported, fixed, and the fix is checked against the",
      "reported case. It passes. Meanwhile the same trick works one street",
      "over, through a door nobody thought was interesting.",
      "This happened in this project, and the second door was found only",
      "because somebody asked 'what else has this shape?'"],
     "After every security fix, hunt the shape rather than the instance — and "
     "write the hunt down (#39, then #40).",
     GREEN),

    ("09_the_mirror.png", "The Mirror", "to answer a question you did not ask",
     ["You ask for spending by region. You are given spending by age group.",
      "It is a real answer, correctly computed, beautifully formatted, and",
      "not what you asked for — so you may well believe something untrue",
      "about the data, and never know."],
     "A request is refused unless the formal query provably answers the question "
     "that was asked; a plausible substitute is not accepted.",
     BLUE),

    ("10_the_imp.png", "The Imp", "the detail nobody counts as an answer",
     ["Numbers are rounded before release, so they cannot be too precise.",
      "But the *order* of the rows was not rounded. Nor was which row got",
      "dropped, nor a statistic computed from the exact count.",
      "Each leaks a little more precision than the rounding allows."],
     "Everything a reader can see must be computed from the published, rounded "
     "numbers alone — checked by a test that perturbs the hidden values and "
     "demands an identical result (#26, #27, #28).",
     GREY),

    ("11_the_parrot.png", "The Parrot", "to repeat hostile text into a trusted place",
     ["The data may contain text somebody else typed. If that text is ever",
      "echoed — into a message on screen, or into the permanent log — then",
      "whoever wrote it has put words into the system's mouth.",
      "The log is meant to be the trustworthy record. That is the target."],
     "Anything echoed back is projected onto a short list of expected words; "
     "anything else is replaced, and the original is stored nowhere (#44).",
     GREEN),

    ("12_the_stampede.png", "The Stampede", "to exhaust what everyone shares",
     ["No single request is an attack. A thousand cheap ones can be —",
      "especially aimed at the one job every other request has to queue",
      "behind, such as verifying the log."],
     "Every route is rate-limited, the shared integrity scan more tightly still, "
     "and each session has a budget of answers (#47).",
     INK),

    ("13_the_doppelganger.png", "The Doppelgänger", "to be two requests at once",
     ["The rule 'have I already answered something too similar?' works by",
      "looking at what has been answered so far. Send both halves of a pair",
      "at the same instant and each looks first, sees nothing, and proceeds.",
      "Two requests, one gap, both allowed."],
     "One person's requests are handled one at a time, across the whole "
     "check-then-record step (#18).",
     BLUE),

    ("14_the_sleeping_dials.png", "The Sleeping Dials",
     "a control switched off and left with its label still reading 'on'",
     ["Not a creature so much as an unlocked cage door with a sign on it.",
      "One safety dial could never actually fire; others would accept settings",
      "that quietly disabled them \u2014 a minimum group size of one, no",
      "rounding \u2014 and still pass every test, because the tests read the",
      "defaults, not the configuration the service was really running."],
     "Floors are enforced on the resolved configuration, the effective policy "
     "is logged at startup, and the only override is a loud environment "
     "variable the config file cannot set for itself (#56, #46).",
     GREY),
]

PACK_HUNTS = [
    ("16_pack_hunt_nixie_rabbit.png", "Two harmless things that are not",
     ["Alone, each of these was a shrug.",
      "",
      "One page listed which exact ages appear in the study — including ages",
      "held by a single person. It never said who.",
      "",
      "Separately, refusals were chatty: the wording told you a little about",
      "why a question could not be answered.",
      "",
      "Together: the first picks the target for free, the second interrogates",
      "them one yes-or-no at a time. Eight questions, every one refused, no",
      "data released — and one person's region, sex, income band and the",
      "type of phone they use, all identified (#29 with #30)."],
     "Lesson: the explanation is an output too. Refusals, error text and even "
     "the word 'denied' spend from the same budget as the data."),

    ("17_pack_hunt_masker_subtractor.png", "Impersonation buys unlimited attempts",
     ["The limits on what one analyst may accumulate — how many answers, and",
      "which combinations are too close together — are counted per name.",
      "",
      "So forging a name did not merely let someone act as a colleague. It",
      "handed them a fresh allowance and an empty history, on demand, as",
      "often as they liked (#45).",
      "",
      "An identity problem and a bookkeeping problem, each modest, combined",
      "into no limits at all."],
     "Lesson: whatever your safety counters are keyed on has just become part "
     "of your identity system, whether you meant it to or not."),

    ("18_pack_hunt_ghost_rabbit.png", "A crash is an answer too",
     ["Every request is meant to leave one line in the tamper-proof log.",
      "",
      "A request that crashed left none \u2014 and whether a given request",
      "crashes depends on the data, so the crash itself is a signal you can",
      "read.",
      "",
      "The unlogged failure and the chatty refusal are one family: a way to",
      "learn something with no released number to account for it (#37)."],
     "Lesson: failure paths are outputs. Audit them, and make a crash "
     "indistinguishable from a data-derived refusal."),

    ("19_pack_hunt_subtractor_hydra.png", "Some cages only hold in pairs",
     ["The round-8 headline attack worked only because two weaknesses lined",
      "up: totals counted rows rather than people, so a one-person difference",
      "hid inside them; and internal range filters could cut between the",
      "published age bands, so a slice could land anywhere.",
      "",
      "Fixing either alone left a working variant. Only closing both, and",
      "reviewing them as one structure, shut the attack (#38 + #39)."],
     "Lesson: some cages are load-bearing pairs. Review them together, which "
     "is exactly what decision record D7 exists to force."),
]


def build(out: str) -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H

    title_slide(prs, "A bestiary of caged beasts",
                "The security problems of a safe data gateway, explained as monsters — "
                "and what each one taught us")

    bullets_slide(prs, "The system, in one slide", [
        "A safe haven holds sensitive records. Nobody may read an individual's data.",
        "Researchers ask questions about GROUPS, and a gateway decides what may leave.",
        "An AI assistant turns English into a formal request — and is NOT trusted:",
        "   it can only tick boxes on a fixed form, and every box is checked again.",
        "This deck is about the ways that arrangement can still go wrong.",
    ], accent=GREEN)

    bullets_slide(prs, "Why monsters?", [
        "The real record is 106 numbered findings, 19 threats and 28 red-team scenarios.",
        "That is complete, precise, and impossible to hold in your head at once —",
        "which matters, because reviews fail when nobody can keep the shape in mind.",
        "People do not remember lists. They remember characters.",
        "Each creature is a memory aid with a real finding number attached to it.",
    ])

    bullets_slide(prs, "How to read a card", [
        "What it wants — the one sentence that tells you when to worry about it.",
        "The story — how it actually got in, in plain English.",
        "The cage — what stops it now, and the finding number, so you can check.",
        "",
        "Nothing here is exterminated. Attacks of this kind are penned, not killed —",
        "and a pen is only as good as whoever notices the door opening.",
    ], accent=GREY)

    for card, name, wants, story, cage, accent in CREATURES:
        card_slide(prs, card, name, wants, story, cage, accent)

    for card, title, story, lesson in PACK_HUNTS:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _bg(slide, WHITE)
        _bar(slide, GREEN)
        art = os.path.join(CARDS, card)
        if os.path.exists(art):
            _fit_centred(slide, art, Inches(0.6), Inches(1.5), Inches(4.4), Inches(5.2))
        _text(slide, Inches(0.6), Inches(0.5), Inches(12.1), Inches(0.8),
              [(title, 30, INK, True)])
        box = _text(slide, Inches(5.4), Inches(1.4), Inches(7.3), Inches(4.2),
                    [(line, 15, INK, False) for line in story])
        for para in box.text_frame.paragraphs:
            para.space_after = Pt(6)
        _text(slide, Inches(5.4), Inches(5.9), Inches(7.3), Inches(1.0),
              [(lesson, 16, GREEN, True)])

    # the meta-specimen gets the full slide; it is the point of the deck
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, INK)
    _bar(slide, BLUE, Inches(0.25))
    art = os.path.join(CARDS, "15_the_blind_zookeeper.png")
    if os.path.exists(art):
        _fit_centred(slide, art, Inches(0.7), Inches(1.2), Inches(4.6), Inches(5.6))
    _text(slide, Inches(5.7), Inches(1.1), Inches(7.0), Inches(1.0),
          [("The Blind Zookeeper", 34, WHITE, True)])
    box = _text(slide, Inches(5.7), Inches(2.2), Inches(7.0), Inches(4.0), [
        ("For seven rounds, this project tested its own defences with a suite "
         "that could not fail.", 18, WHITE, True),
        ("", 8, WHITE, False),
        ("It asked the gateway's own report whether anything had leaked — a "
         "question that, by construction, could only come back 'no'.", 16, WHITE, False),
        ("", 8, WHITE, False),
        ("And it counted any alarm going off anywhere as a pass, which an "
         "attacker can arrange by adding one harmless noisy request.", 16, WHITE, False),
        ("", 8, WHITE, False),
        ("Eleven of the problems in this deck survived that long because the "
         "gaps and the blind test covered for each other.", 16, WHITE, False),
    ])
    for para in box.text_frame.paragraphs:
        para.space_after = Pt(6)
    _text(slide, Inches(5.7), Inches(6.4), Inches(7.0), Inches(0.6),
          [("The test now reads the raw data itself — and is deliberately "
            "given broken defences to check it still says no (#48).", 15,
            BLUE, True)])

    bullets_slide(prs, "Still in the wild — the honest slide", [
        "Not everything is caged, and a map that pretended otherwise would be worse:",
        "",
        "The Straddler — timing still leaks a little; measured, and the number is watched.",
        "The Colluder — two people combining their answers is beyond a per-person limit.",
        "The Residual Head — some combinations are safe in principle, not in practice.",
        "The Optional-Role Imp — a model that omits one table says so by omitting it.",
        "",
        "Each has a measured price and a named fix that would close it properly.",
    ], accent=GREY)

    bullets_slide(prs, "What the beasts taught us", [
        "Name the beast when you cage it — an unnamed fix gets reverted by a refactor.",
        "Every cage needs a keeper, and ask how you know the KEEPER can fail.",
        "Probe the fix, not just the finding — the shape usually outlives the instance.",
        "Test on hostile data, not just hostile questions — refunds and typos suffice.",
        "Price what you cannot close, and say so out loud.",
        "",
        "The full record: docs/bestiary.md, docs/hardening-log.md, docs/specification.md",
    ], accent=GREEN)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, INK)
    _bar(slide, GREEN, Inches(0.25))
    _text(slide, Inches(1.0), Inches(2.9), Inches(11.3), Inches(1.4),
          [("The reserve is never finished.", 34, WHITE, True),
           ("The keepers are the point.", 30, GREEN, True)], align=PP_ALIGN.CENTER)
    _text(slide, Inches(1.0), Inches(6.4), Inches(11.3), Inches(0.6),
          [("Research prototype · synthetic data · not a government service",
            12, GREY, False)], align=PP_ALIGN.CENTER)

    prs.save(out)


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--out", default=os.path.join(ROOT, "artifacts",
                                                  "safe-tre-agent-bestiary.pptx"))
    args = ap.parse_args()
    if not os.path.isdir(CARDS):
        print(f"no cards at {os.path.relpath(CARDS, ROOT)} — run "
              f"scripts/make_bestiary_cards.py first", file=sys.stderr)
        return 1
    os.makedirs(os.path.dirname(args.out), exist_ok=True)
    build(args.out)
    size = os.path.getsize(args.out) / 1e6
    print(f"wrote {os.path.relpath(args.out, ROOT)} ({size:.1f} MB)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
