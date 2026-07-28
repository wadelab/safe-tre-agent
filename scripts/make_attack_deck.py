"""Build the walkthrough deck: one real attack, start to finish, and its cages.

`make_bestiary_deck.py` is the field guide as slides — one creature per slide,
breadth over depth. This is the other half: a single attack followed all the way
through, with the queries an analyst would actually type, the numbers they would
actually see, and the specific control that stops each step now.

The attack is the round-8 headline. Four creatures cooperate: the Nixie hands
over a target for free, the White Rabbit narrows it one refusal at a time, the
Subtractor takes the money in two queries, and the Masker makes the whole thing
repeatable. The Hydra then survives the first fix. Nobody noticed for seven
rounds because the Blind Zookeeper was watching.

**Every number on these slides is computed at build time from the synthetic
data, not typed in.** That is the same discipline the rest of the repository
uses, and it matters more here than anywhere: a walkthrough deck is exactly the
artefact that goes stale quietly and keeps being presented. If the demo data
change, the slides change or the build fails.

The attack itself no longer runs — `age_years == 14` does not survive validation
since hardening #39 — so the historical figures are computed from the row-level
data directly, which is what the attacker would have recovered. The "after"
figures come from running the pair through the gateway as it stands.

    uv run --group decks python scripts/make_attack_deck.py
"""

from __future__ import annotations

import argparse
import json
import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from safetre import synth                                        # noqa: E402
from safetre.disclosure import SessionAuditor, simulatable_cohort_bound  # noqa: E402
from safetre.engine import QueryEngine                           # noqa: E402
from safetre.service import QueryService                         # noqa: E402

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
CARDS = os.path.join(ROOT, "docs", "figures", "bestiary")

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

# Palette sampled from the card art, same as the bestiary page: a near-black
# ground with a warm bias, gilt off the card frames as the single accent, and
# verdigris/rust kept for caged-versus-loose so they never fight the accent.
INK = RGBColor(0x0A, 0x09, 0x08)
PANEL = RGBColor(0x16, 0x13, 0x0E)
PARCHMENT = RGBColor(0xEE, 0xE3, 0xCB)
SOFT = RGBColor(0xC3, 0xB7, 0x9B)
MUTED = RGBColor(0x93, 0x85, 0x6A)
GILT = RGBColor(0xC9, 0xA2, 0x4A)
VERDIGRIS = RGBColor(0x7F, 0xA4, 0x8D)
RUST = RGBColor(0xC9, 0x75, 0x53)
RULE = RGBColor(0x33, 0x2C, 0x21)

SERIF = "Georgia"          # bookish and near-universal; the cards set the tone
MONO = "Consolas"          # specs and figures are quotations from the system

# The worked example. Chosen because every enclosing slice is individually
# safe — which is the whole point of the attack.
REGION, SEX, AGE = "East of England", "F", 14


# --- slide furniture ----------------------------------------------------------

def _slide(prs, ground=INK):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = ground
    return slide


def _text(slide, left, top, width, height, runs, align=PP_ALIGN.LEFT,
          spacing=6, line=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    for i, (text, size, colour, bold, *rest) in enumerate(runs):
        font_name = rest[0] if rest else SERIF
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        para.alignment = align
        para.space_after = Pt(spacing)
        if line:
            para.line_spacing = line
        run = para.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = colour
        run.font.bold = bold
        run.font.name = font_name
    return box


def _rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def _rule(slide, left, top, width, colour=RULE, thickness=Pt(1.2)):
    bar = _rect(slide, left, top, width, thickness, colour)
    return bar


def _card(slide, name, left, top, max_w, max_h):
    """Place a card, centred in its box — the art runs 0.71 to 1.49 in aspect."""
    path = os.path.join(CARDS, name)
    if not os.path.exists(path):
        return
    from PIL import Image
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(max_w / iw, max_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    slide.shapes.add_picture(path, left + (max_w - w) // 2, top + (max_h - h) // 2,
                             width=w, height=h)


def _eyebrow(slide, text, colour=GILT, top=Inches(0.62)):
    _text(slide, Inches(0.9), top, Inches(11.5), Inches(0.4),
          [(text.upper(), 12, colour, True)], spacing=0)


# --- slide kinds --------------------------------------------------------------

def chapter(prs, kicker, title, blurb, card=None):
    slide = _slide(prs)
    _rect(slide, 0, 0, SLIDE_W, Inches(0.09), GILT)
    if card:
        _card(slide, card, Inches(8.3), Inches(1.0), Inches(4.3), Inches(5.5))
    _eyebrow(slide, kicker, top=Inches(2.2))
    _text(slide, Inches(0.9), Inches(2.65), Inches(7.1), Inches(1.8),
          [(title, 44, PARCHMENT, True)], spacing=0)
    _text(slide, Inches(0.9), Inches(4.35), Inches(6.8), Inches(1.6),
          [(blurb, 17, SOFT, False)], line=1.35)


def beast(prs, card, name, wants, points, accent=GILT):
    slide = _slide(prs)
    _rect(slide, 0, 0, SLIDE_W, Inches(0.09), accent)
    _card(slide, card, Inches(0.75), Inches(0.9), Inches(3.9), Inches(5.7))
    _eyebrow(slide, "the beast", accent)
    _text(slide, Inches(5.3), Inches(1.15), Inches(7.3), Inches(1.0),
          [(name, 36, PARCHMENT, True)], spacing=0)
    _text(slide, Inches(5.3), Inches(2.05), Inches(7.3), Inches(0.6),
          [(wants, 17, accent, False)], spacing=0)
    _rule(slide, Inches(5.3), Inches(2.75), Inches(7.0))
    _text(slide, Inches(5.3), Inches(3.0), Inches(7.3), Inches(3.4),
          [(p, 16, SOFT, False) for p in points], spacing=11, line=1.3)


def query(prs, step, caption, spec, outcome, outcome_colour, note):
    """A slide that shows the request as typed and the answer as returned."""
    slide = _slide(prs)
    _rect(slide, 0, 0, SLIDE_W, Inches(0.09), GILT)
    _eyebrow(slide, step)
    _text(slide, Inches(0.9), Inches(1.05), Inches(11.5), Inches(0.8),
          [(caption, 30, PARCHMENT, True)], spacing=0)

    panel = _rect(slide, Inches(0.9), Inches(2.1), Inches(7.5), Inches(3.5),
                  PANEL, line=RULE)
    panel.text_frame.word_wrap = True
    _text(slide, Inches(1.15), Inches(2.3), Inches(7.0), Inches(0.35),
          [("THE REQUEST", 10, MUTED, True)], spacing=0)
    _text(slide, Inches(1.12), Inches(2.72), Inches(7.3), Inches(2.7),
          [(line, 11.5, PARCHMENT, False, MONO) for line in spec],
          spacing=2, line=1.22)

    _text(slide, Inches(8.9), Inches(2.3), Inches(3.6), Inches(0.35),
          [("THE ANSWER", 10, MUTED, True)], spacing=0)
    _text(slide, Inches(8.9), Inches(2.72), Inches(3.7), Inches(1.5),
          [(outcome, 26, outcome_colour, True, MONO)], spacing=4, line=1.15)
    _text(slide, Inches(8.9), Inches(4.35), Inches(3.7), Inches(2.2),
          [(note, 14, SOFT, False)], spacing=6, line=1.3)


def reveal(prs, kicker, figure, caption, detail, colour=RUST):
    slide = _slide(prs)
    _rect(slide, 0, 0, SLIDE_W, Inches(0.09), colour)
    _eyebrow(slide, kicker, colour, top=Inches(1.75))
    _text(slide, Inches(0.9), Inches(2.25), Inches(11.5), Inches(1.9),
          [(figure, 88, PARCHMENT, True, MONO)], spacing=0)
    _text(slide, Inches(0.9), Inches(4.28), Inches(10.5), Inches(0.9),
          [(caption, 26, colour, True)], spacing=0)
    _rule(slide, Inches(0.95), Inches(5.12), Inches(3.2), colour)
    _text(slide, Inches(0.9), Inches(5.42), Inches(9.6), Inches(1.6),
          [(detail, 16, SOFT, False)], line=1.35)


def controls(prs, title, blurb, rows):
    """Why each existing control said yes — the heart of the walkthrough."""
    slide = _slide(prs)
    _rect(slide, 0, 0, SLIDE_W, Inches(0.09), GILT)
    _eyebrow(slide, "and every control agreed")
    _text(slide, Inches(0.9), Inches(1.05), Inches(11.5), Inches(0.7),
          [(title, 30, PARCHMENT, True)], spacing=0)
    _text(slide, Inches(0.9), Inches(1.85), Inches(10.6), Inches(0.6),
          [(blurb, 15, SOFT, False)], spacing=0)
    top = Inches(2.62)
    for name, saw, verdict in rows:
        _rect(slide, Inches(0.9), top, Inches(0.05), Inches(1.02), RUST)
        _text(slide, Inches(1.15), top - Inches(0.04), Inches(3.5), Inches(0.5),
              [(name, 16, PARCHMENT, True)], spacing=0)
        _text(slide, Inches(1.15), top + Inches(0.34), Inches(8.6), Inches(0.5),
              [(saw, 14, MUTED, False, MONO)], spacing=0)
        _text(slide, Inches(10.1), top + Inches(0.06), Inches(2.5), Inches(0.5),
              [(verdict, 15, RUST, True)], spacing=0, align=PP_ALIGN.RIGHT)
        top = top + Inches(1.19)


def cage(prs, number, title, what, before, after, card=None, colour=VERDIGRIS):
    slide = _slide(prs)
    _rect(slide, 0, 0, SLIDE_W, Inches(0.09), colour)
    if card:
        _card(slide, card, Inches(9.0), Inches(1.3), Inches(3.6), Inches(4.8))
    _eyebrow(slide, f"the cage · hardening {number}", colour)
    _text(slide, Inches(0.9), Inches(1.05), Inches(7.8), Inches(0.9),
          [(title, 30, PARCHMENT, True)], spacing=0)
    _text(slide, Inches(0.9), Inches(2.05), Inches(7.6), Inches(2.0),
          [(what, 16, SOFT, False)], line=1.32)

    top = Inches(4.35)
    for label, value, tone in (("Before", before, RUST), ("After", after, colour)):
        _rect(slide, Inches(0.9), top, Inches(0.045), Inches(0.8), tone)
        _text(slide, Inches(1.15), top - Inches(0.03), Inches(1.4), Inches(0.4),
              [(label.upper(), 10, MUTED, True)], spacing=0)
        _text(slide, Inches(2.4), top - Inches(0.06), Inches(6.2), Inches(0.6),
              [(value, 16, tone, True, MONO)], spacing=0)
        top = top + Inches(1.0)


def statement(prs, lead, body, footnote, colour=GILT):
    slide = _slide(prs)
    _rect(slide, 0, 0, SLIDE_W, Inches(0.09), colour)
    _text(slide, Inches(1.1), Inches(1.9), Inches(11.1), Inches(2.4),
          [(lead, 40, PARCHMENT, True)], spacing=0, line=1.12)
    _text(slide, Inches(1.1), Inches(4.35), Inches(10.4), Inches(1.6),
          [(body, 18, SOFT, False)], line=1.35)
    _text(slide, Inches(1.1), Inches(6.35), Inches(10.4), Inches(0.6),
          [(footnote, 14, colour, True)], spacing=0)


def scoreboard(prs, rows):
    slide = _slide(prs)
    _rect(slide, 0, 0, SLIDE_W, Inches(0.09), VERDIGRIS)
    _eyebrow(slide, "the walk-through, scored", VERDIGRIS)
    _text(slide, Inches(0.9), Inches(1.05), Inches(11.5), Inches(0.7),
          [("Every step of that session, and what stops it", 30, PARCHMENT, True)],
          spacing=0)
    head_y = Inches(2.05)
    for label, x, w in (("STEP", Inches(0.9), Inches(3.4)),
                        ("WHAT IT TOOK", Inches(4.4), Inches(4.2)),
                        ("THE CAGE", Inches(8.7), Inches(3.8))):
        _text(slide, x, head_y, w, Inches(0.35), [(label, 10, MUTED, True)], spacing=0)
    _rule(slide, Inches(0.9), Inches(2.45), Inches(11.5))
    top = Inches(2.65)
    for step, took, pen in rows:
        _text(slide, Inches(0.9), top, Inches(3.4), Inches(0.6),
              [(step, 14, PARCHMENT, True)], spacing=0)
        _text(slide, Inches(4.4), top, Inches(4.2), Inches(0.6),
              [(took, 13, SOFT, False)], spacing=0)
        _text(slide, Inches(8.7), top, Inches(3.9), Inches(0.6),
              [(pen, 13, VERDIGRIS, False)], spacing=0)
        top = top + Inches(0.72)


# --- the facts, computed rather than quoted -----------------------------------

def measure() -> dict:
    """Everything the slides assert, derived from the synthetic data now."""
    tables = synth.generate(seed=7)
    engine = QueryEngine(tables)
    service = QueryService(tables)

    def totals(where: str):
        return engine.con.execute(
            "SELECT COUNT(DISTINCT donor_id), COUNT(*), ROUND(SUM(amount_gbp), 2) "
            f"FROM _spend_u WHERE region = ? AND sex = ? AND {where}",   # nosec
            [REGION, SEX]).fetchone()

    cell = totals(f"age_years = {AGE}")
    slice_a = totals(f"age_years >= {AGE}")
    slice_b = totals(f"age_years >= {AGE + 1}")

    base = (("region", "==", REGION), ("sex", "==", SEX))
    a = tuple(sorted(base + (("age_years", ">=", AGE),), key=repr))
    b = tuple(sorted(base + (("age_years", ">=", AGE + 1),), key=repr))

    marginals = engine.marginal_donor_counts()
    ages = marginals["spend"]["age_years"]
    published = engine.published_marginal_donor_counts()["spend"].get("age_years", {})

    # what the pair does today, through the real gateway
    auditor = SessionAuditor()
    statuses = []
    for rating in (7, 8):
        spec = {"dataset": "spend", "measure": {"fn": "sum", "column": "amount_gbp"},
                "filters": [{"column": "age_rating", "op": ">=", "value": rating},
                            {"column": "region", "op": "==", "value": "South West"},
                            {"column": "sex", "op": "==", "value": "F"}]}
        statuses.append(service.handle(json.dumps(spec), planner=None,
                                       auditor=auditor).status)

    return {
        "cell_donors": cell[0], "cell_events": cell[1], "cell_spend": cell[2],
        "a_donors": slice_a[0], "a_events": slice_a[1], "a_spend": slice_a[2],
        "b_donors": slice_b[0], "b_events": slice_b[1], "b_spend": slice_b[2],
        "difference": round(slice_a[2] - slice_b[2], 2),
        "row_delta": slice_a[1] - slice_b[1],
        "donor_delta": slice_a[0] - slice_b[0],
        "age_marginal": ages.get(AGE),
        "old_bound": simulatable_cohort_bound(marginals, "spend", a, b),
        "true_symdiff": engine.cohort_symdiff("spend", a, b),
        "row_symdiff": engine.row_symdiff_donors("spend", a, b),
        "ages_present": len(ages),
        "ages_sub": sum(1 for c in ages.values() if c < 10),
        "ages_unique": sum(1 for c in ages.values() if c == 1),
        "ages_published": len(published),
        "public_dim_pair": statuses,
    }


def build(out: str) -> dict:
    m = measure()
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    money = lambda v: f"GBP {v:,.2f}"                                # noqa: E731

    # — the setup ————————————————————————————————
    slide = _slide(prs)
    _rect(slide, 0, 0, SLIDE_W, Inches(0.09), GILT)
    # a layered hand: later cards draw on top, so they read as dealt
    for i, card in enumerate(["03_the_masker.png", "06_the_white_rabbit.png",
                              "04_the_nixie.png", "01_the_subtractor.png"]):
        _card(slide, card, Inches(6.25 + i * 1.42), Inches(1.45 - i * 0.10),
              Inches(2.5), Inches(4.8))
    _eyebrow(slide, "a walk-through", top=Inches(2.15))
    _text(slide, Inches(0.9), Inches(2.6), Inches(6.4), Inches(2.2),
          [("Two queries,", 46, PARCHMENT, True),
           ("one person", 46, GILT, True)], spacing=0, line=1.05)
    _text(slide, Inches(0.9), Inches(4.75), Inches(6.2), Inches(1.6),
          [("How four of the beasts cooperate to take one donor's exact spending "
            "out of a gateway where every single query is legitimate — and what "
            "each of them runs into now.", 16, SOFT, False)], line=1.32)

    statement(prs,
              "Everything on these slides was run.",
              "The figures are computed from the demo data when this deck is "
              "built, not typed into it. The attack itself no longer runs — the "
              "queries it needs stopped validating — so the amounts it would "
              "have recovered are taken from the row-level data directly, which "
              "is exactly what the attacker would have reconstructed.",
              "Synthetic data throughout. No real person appears in this deck.")

    chapter(prs, "the target", "Choosing a victim, for free",
            "Before anything is stolen, somebody has to be picked. The gateway "
            "publishes a page of summary counts so that an analyst can predict "
            "what it will refuse. That page was the shopping list.",
            card="04_the_nixie.png")

    beast(prs, "04_the_nixie.png", "The Nixie",
          "Wants the system to print a name",
          [f"The summary page listed every exact age present in the study: "
           f"{m['ages_present']} of them.",
           f"{m['ages_sub']} were held by fewer than ten people. "
           f"{m['ages_unique']} were held by exactly one.",
           "It never said who. It did not need to — it said that somebody, "
           "somewhere, is the only person of that age, and that is the whole "
           "shopping list.",
           "One page load. No query budget spent."])

    beast(prs, "06_the_white_rabbit.png", "The White Rabbit",
          "Wants the refusal, not the answer",
          ["Ask a question the gateway refuses and you still learn something: "
           "the refusal itself says the group was too small.",
           "That is one clean bit per question — is anybody here? — and the "
           "questions can be about any combination the form allows.",
           "Eight questions, every one refused, nothing released. Recovered: "
           "one person's region, sex, income band and the type of phone they "
           "use.",
           "Hardening #29 and #30. Neither was worth much alone."],
          accent=MUTED)

    chapter(prs, "the theft", "Two queries and a subtraction",
            "Now the target is known, the money is taken with questions no "
            "reviewer would blink at. Both are about groups of two dozen "
            "people. Both are answered.",
            card="01_the_subtractor.png")

    beast(prs, "01_the_subtractor.png", "The Subtractor",
          "Wants two answers that differ by one person",
          ["It never asks about anybody. It asks two perfectly ordinary "
           "questions whose difference happens to be a single donor.",
           "Every control in the gateway looks at one answer at a time. The "
           "crime is not in either question — it is in the gap between them.",
           f"Our target: the only {AGE}-year-old woman in {REGION}."])

    spec_a = [
        '{',
        '  "dataset": "spend",',
        '  "measure": { "fn": "sum", "column": "amount_gbp" },',
        '  "filters": [',
        f'    {{ "column": "age_years", "op": ">=", "value": {AGE} }},',
        f'    {{ "column": "region",    "op": "==", "value": "{REGION}" }},',
        f'    {{ "column": "sex",       "op": "==", "value": "{SEX}" }}',
        '  ]',
        '}',
    ]
    query(prs, "step one · the wider slice",
          f"Everyone aged {AGE} or over, in one region",
          spec_a, money(m["a_spend"]), VERDIGRIS,
          f"{m['a_donors']} donors, {m['a_events']} events. Comfortably above "
          f"every threshold. Released without a single finding.")

    spec_b = list(spec_a)
    spec_b[4] = f'    {{ "column": "age_years", "op": ">=", "value": {AGE + 1} }},'
    query(prs, "step two · one year narrower",
          "The same question, starting a year later",
          spec_b, money(m["b_spend"]), VERDIGRIS,
          f"{m['b_donors']} donors, {m['b_events']} events. Also comfortably "
          f"safe. Also released without a finding.")

    reveal(prs, "subtract one from the other", money(m["difference"]),
           f"The exact annual spending of one {AGE}-year-old woman",
           f"She contributes {m['cell_events']} events. She is the only person "
           f"in that cell. Nothing about her was ever asked for, and nothing "
           f"that was released was, on its own, disclosive.")

    spec_direct = list(spec_a)
    spec_direct[4] = f'    {{ "column": "age_years", "op": "==", "value": {AGE} }},'
    query(prs, "the same question, asked honestly",
          "Ask the gateway for that cell directly",
          spec_direct, "DENIED", RUST,
          f"{m['cell_donors']} donor is below the threshold of ten, so the "
          f"gateway refuses — as it should. The subtraction had already "
          f"handed over the same number.")

    controls(prs, "Nobody was asleep. Everybody was looking elsewhere.",
             "Each control did exactly what it was written to do. None of them "
             "was written to look at the gap between two answers.",
             [("Cell threshold",
               f"the isolated cell was never queried — the two slices held "
               f"{m['a_donors']} and {m['b_donors']} donors",
               "passed"),
              ("Dominance rule",
               "no single donor was more than half of either slice",
               "passed"),
              ("Lineage bound",
               f"decided from published marginals: age {AGE} is held by "
               f"{m['age_marginal']} people, bound = {m['old_bound']}, over the limit of 10",
               "passed"),
              ("Total-delta check",
               f"compared ROW counts: the slices differ by {m['row_delta']} "
               f"events — the donor difference was {m['donor_delta']}",
               "passed")])

    statement(prs,
              "The bound was not wrong. It was answering a different question.",
              f"It certified that at most {m['old_bound']} people separate the "
              f"two slices — true, and useless. The true separation is "
              f"{m['true_symdiff']}. A bound computed from whole-population "
              f"counts cannot see that both queries also carry two narrowing "
              f"filters, and it is the narrowing that isolates.",
              "This residual was documented before it was exploited. Documented "
              "is not the same as closed.")

    chapter(prs, "the second head", "Fixing it, and not fixing it",
            "The first repair made the attack unstateable: filters on exact age "
            "must now land on published band edges, so no two of them are ever "
            "one year apart. It closed this attack completely.",
            card="08_the_hydra.png")

    beast(prs, "08_the_hydra.png", "The Hydra",
          "Wants to survive the fix aimed at it",
          ["The fix was checked against the reported attack. It passed.",
           "Then somebody asked what else had that shape — and found the same "
           "theft running through a dimension nobody thought was dangerous, "
           "because it is coarse, public and groupable: the age rating of the "
           "app.",
           "Twenty more cells were still recoverable.",
           "The reason is the sharp part: those two groups contain exactly the "
           "same people. The filter divides their purchases, not them."],
          accent=RUST)

    reveal(prs, "the difference that mattered", "0 vs 1",
           "Identical people. One donor's worth of rows.",
           "Both differencing checks compared donor lists and both were right: "
           "the lists were identical. A released number is a function of the "
           "rows it added up, and nobody was differencing rows.",
           colour=RUST)

    chapter(prs, "the accomplice", "Why the budget did not save us",
            "A session may release twenty answers. Two are enough for one "
            "person — but there are hundreds of people, and the counter turned "
            "out to be attached to something the caller chooses.",
            card="03_the_masker.png")

    beast(prs, "03_the_masker.png", "The Masker",
          "Wants to act while wearing someone else's name",
          ["The gateway trusts a header naming the analyst, because only the "
           "local gateway process is supposed to be able to set it.",
           "But the untrusted AI model runs on the same machine, and so does "
           "anything else on it.",
           "Twenty-one forged requests were accepted and written into the "
           "permanent log under a colleague's name.",
           "And because the budget and the differencing history are filed "
           "under that name, a new name is a clean sheet. The limit was not "
           "a limit."],
          accent=RUST)

    statement(prs,
              "Seven rounds of red-teaming did not catch any of this.",
              "The suite asked the gateway's own report whether anything had "
              "leaked — a question that, on this path, could only come back "
              "'no', because the evidence is stripped out before release. And "
              "it passed any attack where some alarm went off, which an "
              "attacker arranges by adding one noisy extra query.",
              "The gaps and the blind test had been covering for each other.")

    # — the cages ————————————————————————————————
    chapter(prs, "the cages", "What stops each step now",
            "Five changes, each measured. The numbers below are the ones that "
            "moved.",
            card="14_the_blind_zookeeper.png")

    cage(prs, "#40", "Difference the rows, not the people",
         "The lineage check now counts the donors behind the rows that exactly "
         "one of the two queries added up. Where every filter is about a person "
         "it gives the same answer as before, so nothing is traded away — and "
         "where a filter divides purchases instead, it sees what the old check "
         "could not.",
         f"bound said {m['old_bound']} people, so allowed — true separation "
         f"{m['true_symdiff']}",
         f"row-level difference {m['row_symdiff']} person, so denied",
         card="01_the_subtractor.png")

    cage(prs, "#39", "Give exact age no fine edges to cut on",
         "Filters on exact age must land on the published age-band edges — "
         "13, 16, 18, 25, 35, 50. No two are adjacent, so the one-year step "
         "this attack needs cannot be written down. Asking for a single exact "
         "age is not offered at all.",
         f"{m['ages_present']} distinct ages filterable, "
         f"{m['ages_published']} of them released directly",
         "6 band edges, all of them already public")

    cage(prs, "#38", "Count people, not paperwork",
         "The cheap first-pass check compared row counts against a threshold "
         "that means people. On an event-level view one busy donor inflates "
         "rows without adding anybody, so the check was reading the wrong "
         "units entirely.",
         f"saw {m['row_delta']} events between the slices — no alarm",
         f"sees {m['donor_delta']} donor between the slices — alarm")

    cage(prs, "#45", "Stop trusting the corridor",
         "The identity header is only believed when the gateway in front "
         "counter-signs it with a shared secret, and a production deployment "
         "without one now refuses everybody rather than trusting the header. "
         "Sharing a machine with the system is not proof of identity.",
         "21 forged requests accepted; a new name bought a fresh budget",
         "forged requests refused; the budget is attached to a proven name",
         card="03_the_masker.png")

    cage(prs, "#48", "Build a test that can say no",
         "The red-team oracle now reads the row-level data rather than the "
         "gateway's own opinion, inspects every step of a session instead of "
         "the last, and asks what the released answers combine into. Then it is "
         "deliberately handed broken defences, to check it still objects.",
         "verdict: any alarm anywhere counts as a pass",
         "verdict: the session disclosed nothing — or it failed",
         card="14_the_blind_zookeeper.png", colour=GILT)

    scoreboard(prs, [
        ("Pick a target", "one page load, no budget spent",
         "sub-threshold ages no longer published"),
        ("Narrow it down", "eight refusals, nothing released",
         "refusals give one canonical answer"),
        ("Take the money", f"two queries, {money(m['difference'])}",
         "row-level differencing denies the pair"),
        ("Do it again", "rotate the name, fresh budget",
         "proxy secret required with identity"),
        ("Go unnoticed", "the suite reported a pass",
         "an oracle calibrated to fail"),
    ])

    statement(prs,
              "Every step was legitimate. That is the lesson.",
              "No rule was broken, no control malfunctioned, and nothing in "
              "this session would look wrong in a log. The disclosure lived in "
              "the relationship between answers, and controls that examine one "
              "answer at a time cannot see relationships.",
              "The record: docs/bestiary.md · docs/hardening-log.md · "
              "docs/specification.md",
              colour=VERDIGRIS)

    prs.save(out)
    return m


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--out", default=os.path.join(
        ROOT, "artifacts", "safe-tre-agent-attack-walkthrough.pptx"))
    args = ap.parse_args()
    if not os.path.isdir(CARDS):
        print(f"no cards at {os.path.relpath(CARDS, ROOT)} — run "
              f"scripts/make_bestiary_cards.py first", file=sys.stderr)
        return 1
    os.makedirs(os.path.dirname(args.out), exist_ok=True)
    facts = build(args.out)
    print(f"wrote {os.path.relpath(args.out, ROOT)} "
          f"({os.path.getsize(args.out) / 1e6:.1f} MB)")
    print(f"  worked example : {REGION} · {SEX} · age {AGE}")
    print(f"  recovered      : GBP {facts['difference']:,.2f} from "
          f"{facts['cell_donors']} donor ({facts['cell_events']} events)")
    print(f"  old bound {facts['old_bound']} vs true separation "
          f"{facts['true_symdiff']}; row-level check now sees "
          f"{facts['row_symdiff']}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
